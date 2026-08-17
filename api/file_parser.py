"""
File parser module for StudentAI.
Extracts text content from PDF, PowerPoint, and Word documents.
"""

import os
import tempfile
from typing import Optional


def extract_text_from_pdf(file_path: str) -> str:
    """Extract text from a PDF file using PyMuPDF."""
    try:
        import fitz  # PyMuPDF
        text_parts = []
        doc = fitz.open(file_path)
        for page_num in range(len(doc)):
            page = doc[page_num]
            text = page.get_text("text")
            if text.strip():
                text_parts.append(f"--- Page {page_num + 1} ---\n{text}")
        doc.close()
        return "\n\n".join(text_parts)
    except Exception as e:
        print(f"Error extracting PDF text: {e}")
        return ""


def extract_text_from_pptx(file_path: str) -> str:
    """Extract text from a PowerPoint file."""
    try:
        from pptx import Presentation
        text_parts = []
        prs = Presentation(file_path)
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
            if slide_text:
                text_parts.append(f"--- Slide {slide_num} ---\n" + "\n".join(slide_text))
        return "\n\n".join(text_parts)
    except Exception as e:
        print(f"Error extracting PPTX text: {e}")
        return ""


def extract_text_from_docx(file_path: str) -> str:
    """Extract text from a Word document."""
    try:
        from docx import Document
        text_parts = []
        doc = Document(file_path)
        for para in doc.paragraphs:
            if para.text.strip():
                text_parts.append(para.text)
        # Also extract from tables
        for table in doc.tables:
            for row in table.rows:
                row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if row_text:
                    text_parts.append(" | ".join(row_text))
        return "\n\n".join(text_parts)
    except Exception as e:
        print(f"Error extracting DOCX text: {e}")
        return ""


def extract_text_from_txt(file_path: str) -> str:
    """Extract text from a plain text file."""
    try:
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()
    except Exception as e:
        print(f"Error reading TXT file: {e}")
        return ""


def extract_text(file_path: str, file_type: str) -> str:
    """
    Extract text from a file based on its type.
    
    Args:
        file_path: Path to the uploaded file
        file_type: File extension (pdf, pptx, docx, txt)
    
    Returns:
        Extracted text content as a string
    """
    file_type = file_type.lower().strip(".")
    
    extractors = {
        "pdf": extract_text_from_pdf,
        "pptx": extract_text_from_pptx,
        "ppt": extract_text_from_pptx,
        "docx": extract_text_from_docx,
        "doc": extract_text_from_docx,
        "txt": extract_text_from_txt,
        "text": extract_text_from_txt,
    }
    
    extractor = extractors.get(file_type)
    if extractor is None:
        print(f"Unsupported file type: {file_type}")
        return ""
    
    text = extractor(file_path)
    
    # Clean up the text
    if text:
        # Remove excessive whitespace
        lines = text.split("\n")
        cleaned_lines = [line.strip() for line in lines]
        text = "\n".join(line for line in cleaned_lines if line)
    
    return text


def get_file_type(filename: str) -> Optional[str]:
    """Get the file type from a filename."""
    if "." in filename:
        return filename.rsplit(".", 1)[1].lower()
    return None
